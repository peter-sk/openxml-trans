from pptx import Presentation

from .base_trans import BaseTranslator

class PPTXTranslator(BaseTranslator):

    def __init__(self, model_name=None, service_url=None, cache_name=None, device=None, clean_text=False, original_language=None, result_language=None):
        super(self.__class__, self).__init__(model_name=model_name, service_url=service_url, cache_name=cache_name, device=device, clean_text=clean_text, original_language=original_language, result_language=result_language)

    def runs_compatible(self, r1, r2):
        return r1.font.size == r2.font.size and r1.font.name == r2.font.name and r1.font.bold == r2.font.bold and r1.font.italic == r2.font.italic

    def add_run(self, para, run, new_text):
        new_run = para.add_run()
        new_run.text = new_text
        new_run.font.size = run.font.size
        new_run.font.name = run.font.name
        new_run.font.bold = run.font.bold
        new_run.font.italic = run.font.italic

    def translate_document(self, document):
        for slide in document.slides:
            for shape in slide.shapes:
                for para in shape.text_frame.paragraphs:
                    self.translate_para(para)

    def load_document(self, document_name):
        return Presentation(document_name)

    def save_document(self, document, document_name):
        document.save(document_name)
