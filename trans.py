from cleantext import clean
from docx import Document
from os.path import isfile
from pickle import load, dump
from pptx import Presentation
from sys import argv
from torch.cuda import is_available
from transformers import pipeline
device = "cuda:0" if is_available() else "cpu"
p = pipeline("translation", "Helsinki-NLP/opus-mt-da-en", device=device)
if isfile("cache"):
    cache = load(open("cache", "rb"))
else:
    cache = {}
def translate(text):
    if text in cache:
        return cache[text]
    new_sentences = []
    sentences = clean(text, lang="da", lower=False).split(". ")
    for i in range(len(sentences)):
        sentence = sentences[i]
        if i+1 < len(sentences):
            sentence += "."
        new_sentence = p(sentence)[0]['translation_text'] if any((c.isalpha() for c in sentence)) else sentence
        print("REPR",repr(sentence), repr(new_sentence))
        new_sentences.append(new_sentence)
    new_text = " ".join(new_sentences)
    cache[text] = new_text
    dump(cache, open("cache", "wb"))
    return new_text
def compatible(r1, r2):
    if hasattr(r1, "style"):
        return r1.font.size == r2.font.size and r1.font.name == r2.font.name and r1.bold == r2.bold and r1.italic == r2.italic and r1.style == r2.style
    else:
        return r1.font.size == r2.font.size and r1.font.name == r2.font.name and r1.font.bold == r2.font.bold and r1.font.italic == r2.font.italic
def translate_para(para):
    runs = list(para.runs)
    print("OLD","|".join((run.text for run in runs)))
    new_runs = []
    while runs:
        run, runs = runs[0], runs[1:]
        text = run.text
        while runs and compatible(run, runs[0]):
            text += runs[0].text
            runs = runs[1:]
        run.text = text
        new_runs.append(run)
    para.clear()
    for run in new_runs:
        new_text = translate(run.text)
        if hasattr(run, "style"):
            para.add_run(text=new_text, style=run.style)
            new_run = para.runs[-1]
            new_run.font.size = run.font.size
            new_run.font.name = run.font.name
            new_run.bold = run.bold
            new_run.italic = run.italic
        else:
            new_run = para.add_run()
            new_run.text = new_text
            new_run.font.size = run.font.size
            new_run.font.name = run.font.name
            new_run.font.bold = run.font.bold
            new_run.font.italic = run.font.italic
    print("NEW","|".join((run.text for run in para.runs)))

if __name__ == "__main__":
    ext = argv[1].split(".")[-1]
    if ext == "docx":
        print("DOCX")
        d = Document(argv[1])
        for para in d.paragraphs:
            translate_para(para)
        for table in d.tables:
            for row in table.rows:
                for cell in row.cells:
                    for para in cell.paragraphs:
                        translate_para(para)
        d.save(argv[2])
    elif ext == "pptx":
        print("PPTX")
        d = Presentation(argv[1])
        for slide in d.slides:
            for shape in slide.shapes:
                for para in shape.text_frame.paragraphs:
                    translate_para(para)
        d.save(argv[2])
    else:
        print(f"UNKNOWN extension {ext}")